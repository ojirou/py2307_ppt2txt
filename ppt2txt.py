import os
import glob
import pptx
from pptx.util import Pt
def extract_text_from_pptx(file_path):
    presentation = pptx.Presentation(file_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text.strip() + " "
                    text += "\n"
    return text
def remove_empty_lines(file_path):
    with open(file_path, "r+", encoding="utf-8") as file:
        lines = file.readlines()
        lines = [line.strip() for line in lines if line.strip()]
        file.seek(0)
        file.truncate()
        file.write("\n".join(lines))
def process_pptx_files(folder_path):
    pptx_files = glob.glob(os.path.join(folder_path, "*.pptx"))
    for file_path in pptx_files:
        # Extract text from pptx
        text = extract_text_from_pptx(file_path)
        # Generate output text file path
        file_name = os.path.splitext(os.path.basename(file_path))[0]
        text_file_path = os.path.join(folder_path, file_name + ".txt")
        # Write text to file
        with open(text_file_path, "w", encoding="utf-8") as text_file:
            text_file.write(text)
        # Remove empty lines from text file
        remove_empty_lines(text_file_path)
folder_path = input("指定フォルダのパスを入力してください: ")
# 入力されたフォルダの存在を確認
if not os.path.exists(folder_path):
    print("指定されたフォルダが存在しません。")
    return
# 処理を実行
process_pptx_files(folder_path)