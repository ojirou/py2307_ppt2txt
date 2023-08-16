import os
import subprocess
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
# from pptx.util import Presentation
# from pptx.util import Presentation
def check_recursively_for_text(this_set_of_shapes, txt_list):
    for shape in this_set_of_shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            check_recursively_for_text(shape.shapes, txt_list)
        else:
            if hasattr(shape, "text"):
                txt_list.append(shape.text)
            if shape.has_table:
                for cell in shape.table.iter_cells():
                    txt_list.append(cell.text)
    return txt_list
if __name__=="__main__":
    folder_path = r"C:\Users\user\python\IMAGE\ppt_file"
    # output_folder_path = r"C:\Users\user\python\PDF\text_output"
    # os.makedirs(output_folder_path, exist_ok=True)
    for file_name in os.listdir(folder_path):
        if file_name.endswith(".pptx"):
            prs = Presentation(os.path.join(folder_path, file_name))
            all_slide_text=[]
            for i, slide in enumerate(prs.slides):
                # extract_text_from_pdf(file_path)
                txt_list=[]
                txt_list=check_recursively_for_text(slide.shapes, txt_list)
                if slide.has_notes_slide and slide.notes_slides.notes_text_frame.text:
                    txt_list.append(slide.notes_slide.notes_text_frame.text)
                txt_list=[line for line in txt_list if line.strip()]
                all_slide_text.append('\n'.join(txt_list))
            slide_name=os.path.splitext(file_name)[0]
            folder_name=r"C:\Users\user\python\ppt\slide_txt_output"
            os.makedirs(folder_name, exist_ok=True)
            output_file_path=os.path.join(folder_name, f"output_{slide_name}.txt")
            with open(output_file_path, "w", encoding="utf-8") as f:
                 f.write('\n\n'.join(all_slide_text))
            subprocess.Popen(["notepad.exe", output_file_path])